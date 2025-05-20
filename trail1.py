import os
import re
import streamlit as st
import pyttsx3
import google.generativeai as genai  # Import Gemini
from pptx import Presentation
from io import BytesIO
from langchain_community.document_loaders import PyPDFLoader
from dotenv import load_dotenv
import os
import tempfile
import streamlit as st
from gtts import gTTS
load_dotenv()

# Get the API key from environment variables
api_key = st.secrets["GEMINI_API_KEY"]
# Initialize Gemini 1.5 Flash
genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-1.5-flash")

# Streamlit UI
st.title("Research Paper Summarizer and PowerPoint Generator with 🎙 Podcast Feature")

st.write("Upload a PDF file to get a summary, listen to expert discussion audio, and download the PowerPoint presentation.")

# Dropdown for summary level
level_prompts = {
    "Beginner": "Summarize this research paper section for a high school student in 4-6 concise bullet points:",
    "Student": "Create a structured summary of this section for undergraduate students in 4-6 points:",
    "Expert": "Generate a detailed summary of this section for researchers in 5-7 well-formed bullet points:"
}
summary_level = st.selectbox("Select summary level:", list(level_prompts.keys()))

# Dropdown for creativity level
creativity_levels = {
    "Formal": "Keep the conversation strictly professional and formal.",
    "Balanced": "Maintain a balance between professional and conversational tone.",
    "Creative": "Make the conversation more creative and engaging with some informal elements."
}
creativity_level = st.selectbox("Select creativity level:", list(creativity_levels.keys()))

# Dropdown for podcast length
podcast_lengths = {
    "Short (2-3 mins)": "Generate a short podcast with 2-3 questions and concise answers.",
    "Medium (5-7 mins)": "Generate a medium-length podcast with 4-5 questions and detailed answers.",
    "Long (10+ mins)": "Generate a long podcast with 6-8 questions and in-depth discussion."
}
podcast_length = st.selectbox("Select podcast length:", list(podcast_lengths.keys()))

# Dropdown for PowerPoint template selection
template_options = {
    "Template 1": "theme_template_1.pptx",
    "Template 2": "theme_template_2.pptx",
    "Template 3": "theme_template_3.pptx"
}
selected_template = st.selectbox("Select PowerPoint Template:", list(template_options.keys()))

# Custom CSS for avatars
st.markdown("""
<style>
    .avatar-container {
        display: flex;
        justify-content: space-around;
        margin: 2rem 0;
        padding: 20px;
        background: #f0f2f6;
        border-radius: 15px;
    }
    .avatar-card {
        text-align: center;
        transition: all 0.3s ease;
        padding: 15px;
        border-radius: 10px;
    }
    .avatar-img {
        width: 120px;
        height: 120px;
        border-radius: 50%;
        transition: all 0.3s ease;
        filter: grayscale(80%) brightness(0.8);
        border: 3px solid transparent;
    }
    .active-speaker .avatar-img {
        filter: none;
        transform: scale(1.1);
        border-color: #4CAF50;
        box-shadow: 0 0 20px rgba(76, 175, 80, 0.3);
    }
    .avatar-name {
        margin-top: 10px;
        font-weight: bold;
        color: #666;
    }
    .active-speaker .avatar-name {
        color: #4CAF50;
    }
    @keyframes pulse {
        0% { transform: scale(1); }
        50% { transform: scale(1.05); }
        100% { transform: scale(1); }
    }
    .speaking-indicator {
        width: 15px;
        height: 15px;
        background: #4CAF50;
        border-radius: 50%;
        margin: 10px auto;
        opacity: 0;
        animation: pulse 1s infinite;
    }
    .active-speaker .speaking-indicator {
        opacity: 1;
    }
</style>
""", unsafe_allow_html=True)

# File uploader for PDF
pdf_file = st.file_uploader("Upload PDF", type=["pdf"])

def get_avatar_html(active_speaker):
    return f"""
    <div class="avatar-container">
        <div class="avatar-card {'active-speaker' if active_speaker == 'Alex' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-female-skin-type-5.png" 
                 class="avatar-img">
            <div class="avatar-name">Alex</div>
        </div>
        <div class="avatar-card {'active-speaker' if active_speaker == 'Dr. Smith' else ''}">
            <div class="speaking-indicator"></div>
            <img src="https://img.icons8.com/color/144/000000/circled-user-male-skin-type-7.png" 
                 class="avatar-img">
            <div class="avatar-name">Dr. Smith</div>
        </div>
    </div>
    """

def extract_and_summarize_sections(text, summary_level):
    """Extract sections and generate summaries using Gemini."""
    prompt = f"""Analyze the following research paper and:
1. Identify all major sections.
2. For each section, generate a summary using the following guidelines:
   - {level_prompts[summary_level]}
3. Format the response as:
   ## Section Name
   - Bullet point 1
   - Bullet point 2
   - Bullet point 3

Paper content:
{text}
"""
    try:
        response = model.generate_content(prompt)  # Use the initialized model
        return response.text
    except Exception as e:
        st.error(f"Error processing document: {e}")
        return None

def create_ppt_from_summary(summary_text, template_path):
    """Create PowerPoint from section-wise summaries using the selected template."""
    prs = Presentation(template_path)

    # Generate a title for the presentation using Gemini
    title_prompt = f"""Analyze the following text and generate a concise, professional title for a PowerPoint presentation (maximum 10-12 words):
    {summary_text[:5000]}  # Use the first 5000 characters for title generation
    """
    try:
        title_response = model.generate_content(title_prompt)
        title = title_response.text.strip()
    except Exception as e:
        st.error(f"Error generating title: {e}")
        title = "Research Summary"  # Default title if title generation fails

    # Ensure the title fits within one or two lines
    max_title_length = 80  # Maximum characters per line (adjust as needed)
    if len(title) > max_title_length:
        # Split the title into two lines
        words = title.split()
        line1 = ""
        line2 = ""
        for word in words:
            if len(line1) + len(word) + 1 <= max_title_length:  # +1 for space
                line1 += word + " "
            else:
                line2 += word + " "
        title = f"{line1.strip()}\n{line2.strip()}"

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]  # 0 is the layout for a title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated by AI Podcast Generator"

    # Add section slides
    slide_layout = prs.slide_layouts[1]  # 1 is the layout for a content slide

    # Parse the summary text into sections and bullet points
    sections = {}
    current_section = "Introduction"
    for line in summary_text.split('\n'):
        if line.startswith("## "):
            current_section = line[3:].strip()
            sections[current_section] = []
        elif line.startswith("- "):
            sections[current_section].append(line[2:])

    # Create slides for each section
    for section, bullets in sections.items():
        slides_per_section = min((len(bullets) // 6) + 1, 5)
        chunk_size = max(len(bullets) // slides_per_section, 1)
        
        for i in range(0, len(bullets), chunk_size):
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section if i == 0 else f"{section} (Cont.)"
            
            content_box = slide.shapes.placeholders[1]
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for bullet in bullets[i:i+chunk_size]:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                p.space_after = 0
    
    pptx_stream = BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

def generate_podcast_script(summary_text, creativity_level, podcast_length):
    """Generate a conversational podcast script using Gemini."""
    prompt = f"""Create a conversational podcast script between host Alex and researcher Dr. Smith discussing the research paper. Follow these rules:
1. Alex should ask curious, layperson-friendly questions.
2. Dr. Smith should provide expert answers based on the paper.
3. Always prefix lines with either "Alex:" or "Dr. Smith:".
4. Keep responses conversational but informative.
5. Cover key findings, methodology, and implications.
6. {creativity_levels[creativity_level]}
7. {podcast_lengths[podcast_length]}

Paper summary:
{summary_text}
"""
    try:
        response = model.generate_content(prompt)  # Use the initialized model
        return response.text
    except Exception as e:
        st.error(f"Error generating podcast script: {e}")
        return None

def generate_podcast_audio(podcast_script, rate=150):
    """Generate TTS audio with distinct voices for host and researcher in conversation order."""
    try:
        # Create a temporary directory to store the audio files
        temp_dir = tempfile.mkdtemp()
        
        # Parse script into segments
        segments = []
        for line in podcast_script.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith("Alex:"):
                speaker = "Alex"
                text = line.replace("Alex:", "").strip()
                segments.append((speaker, text))
            elif line.startswith("Dr. Smith:"):
                speaker = "Dr. Smith"
                text = line.replace("Dr. Smith:", "").strip()
                segments.append((speaker, text))
        
        # Generate individual audio files for each segment
        audio_files = []
        for i, (speaker, text) in enumerate(segments):
            segment_file = os.path.join(temp_dir, f"segment_{i}.mp3")
            
            # Use different TTS settings based on speaker
            if speaker == "Alex":
                tts = gTTS(text=text, lang='en', tld='com.au', slow=False)  # Australian English for female host
            else:  # Dr. Smith
                tts = gTTS(text=text, lang='en', tld='co.uk', slow=False)  # British English for male researcher
            
            tts.save(segment_file)
            audio_files.append((speaker, segment_file))
        
        # Display the complete conversation audio
        st.write("### 🎧 Complete Podcast Audio")
        
        # Create navigation for audio segments
        if 'current_segment' not in st.session_state:
            st.session_state.current_segment = 0
        
        # Avatar display area
        avatar_placeholder = st.empty()
        
        # Show current speaker's avatar
        if st.session_state.current_segment < len(segments):
            current_speaker = segments[st.session_state.current_segment][0]
            avatar_placeholder.markdown(get_avatar_html(current_speaker), unsafe_allow_html=True)
        else:
            avatar_placeholder.markdown(get_avatar_html(''), unsafe_allow_html=True)
        
        # Audio player for current segment
        if st.session_state.current_segment < len(audio_files):
            speaker, audio_file = audio_files[st.session_state.current_segment]
            with open(audio_file, "rb") as f:
                audio_bytes = f.read()
                st.audio(audio_bytes, format="audio/mp3")
                st.markdown(f"**Now Speaking: {speaker}**")
        
        # Navigation controls
        col1, col2, col3, col4 = st.columns([1, 1, 1, 1])
        
        with col1:
            if st.button("⏮️ Reset"):
                st.session_state.current_segment = 0
                st.experimental_rerun()
                
        with col2:
            if st.button("⏪ Previous") and st.session_state.current_segment > 0:
                st.session_state.current_segment -= 1
                st.experimental_rerun()
                
        with col3:
            if st.button("Next ⏩") and st.session_state.current_segment < len(segments) - 1:
                st.session_state.current_segment += 1
                st.experimental_rerun()
                
        with col4:
            if st.button("Auto-Play All"):
                # In Streamlit, we can't truly auto-play sequentially
                # This is a limitation of the Streamlit framework
                st.warning("Due to Streamlit limitations, we can't auto-play segments sequentially. You can manually click through the segments.")
        
        # Add option to download full conversation
        if st.button("Generate & Download Full Conversation"):
            # This would merge all audio files into one
            from pydub import AudioSegment
            
            try:
                # Use pydub to concatenate audio files
                combined = AudioSegment.empty()
                for _, file_path in audio_files:
                    segment = AudioSegment.from_mp3(file_path)
                    combined += segment
                
                # Save the combined audio
                full_audio_path = os.path.join(temp_dir, "full_conversation.mp3")
                combined.export(full_audio_path, format="mp3")
                
                # Provide download button
                with open(full_audio_path, "rb") as f:
                    full_audio_bytes = f.read()
                    st.download_button(
                        label="Download Full Conversation Audio",
                        data=full_audio_bytes,
                        file_name="podcast_conversation.mp3",
                        mime="audio/mp3"
                    )
            except Exception as e:
                st.error(f"Error generating full audio: {e}")
                st.info("Note: You may need to install the pydub package and ffmpeg. Add 'pydub' to your requirements.txt file.")
        
        # Display script with current line highlighted
        st.write("### 📝 Podcast Script")
        for idx, (speaker, text) in enumerate(segments):
            if idx == st.session_state.current_segment:
                # Highlight current segment
                st.markdown(f"""
                <div style='background-color: #e6f7ff; padding: 10px; border-radius: 5px; border-left: 5px solid #1890ff;'>
                <strong>{speaker}:</strong> {text}
                </div>
                """, unsafe_allow_html=True)
            else:
                # Normal display for other lines
                if speaker == "Alex":
                    st.markdown(f"<div style='color:#FF5733; padding:5px;'><strong>{speaker}:</strong> {text}</div>", unsafe_allow_html=True)
                else:
                    st.markdown(f"<div style='color:#3366FF; padding:5px;'><strong>{speaker}:</strong> {text}</div>", unsafe_allow_html=True)
        
        # Clean up temporary files when done
        for _, file_path in audio_files:
            if os.path.exists(file_path):
                os.remove(file_path)
        
        # Remove the full conversation file if it exists
        full_audio_path = os.path.join(temp_dir, "full_conversation.mp3")
        if os.path.exists(full_audio_path):
            os.remove(full_audio_path)
            
        os.rmdir(temp_dir)
        
        return True
        
    except Exception as e:
        st.error(f"Error generating audio: {e}")
        return False


if pdf_file is not None:
    # Save the uploaded PDF temporarily
    temp_pdf_path = f"./temp_{pdf_file.name}"
    with open(temp_pdf_path, "wb") as f:
        f.write(pdf_file.getbuffer())

    # Load and process the PDF
    loader = PyPDFLoader(temp_pdf_path)
    documents = loader.load()
    pdf_text = "\n".join([doc.page_content for doc in documents])

    # Extract sections and generate summaries using Gemini
    with st.spinner("Analyzing document and generating summaries..."):
        summary_text = extract_and_summarize_sections(pdf_text, summary_level)

    if summary_text:
        # Display summaries
        st.subheader("Section-wise Summary")
        st.markdown(summary_text)

        # Generate podcast content
        podcast_script = generate_podcast_script(summary_text, creativity_level, podcast_length)

        if podcast_script:
            st.subheader("Expert Discussion Script")
            st.text_area(
            "Expert Discussion Script",
            value=podcast_script,
            height=300,  # Set the height of the textarea
            key="podcast_script_area"
             )
            
            # Play podcast audio
            st.subheader("Listen to Expert Discussion")
            if st.button("Play Podcast"):
                if generate_podcast_audio(podcast_script):
                    st.success("Podcast audio played successfully!")
                else:
                    st.error("Failed to play podcast audio.")

        # Generate and download PowerPoint
        pptx_stream = create_ppt_from_summary(summary_text, template_options[selected_template])
        
        st.subheader("Download Presentation Slides")
        st.download_button(
            label="Download PowerPoint",
            data=pptx_stream,
            file_name="research_summary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    # Cleanup
    os.remove(temp_pdf_path)
